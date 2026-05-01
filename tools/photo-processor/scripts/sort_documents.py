#!/usr/bin/env python3
"""
Sort Documents & Media — organizes documents, video, audio, design files, and archives.

Handles: PDF, Word, Excel, PowerPoint, CSV, text, Pages/Numbers/Keynote, ePub,
         MOV, MP4, AVI, MKV, MP3, WAV, FLAC, M4A, PSD, AI, Sketch, ZIP, DMG, etc.

Grouping logic (priority order):
  0. Source folder name            — Your Company/ folder → Projects/Sample_Co/ (or Companies/)
  1. Name-based project grouping   — files sharing a keyword → Projects/<Name>/
  2. Vault entity match            — known contact/company → People/ or Companies/
  3. Content-based category (docs) — Finance, Legal, Medical, Work, etc.
  4. Fallback                      — file type + month

Output structure:
  Organized_Docs/
  ├── People/Charlie_Smith/           ← doc mentions a known contact
  ├── Companies/Sample_Co/             ← matches vault service + has .mov + .pdf
  ├── Projects/YOURBRAND/                  ← source folder or shared filename keyword
  │   ├── brand_guidelines.pdf
  │   ├── promo_cut_v2.mov
  │   └── logo_final.psd
  ├── Finance/2025-01/
  ├── Legal/2025-03/
  ├── Video/2025-07/                  ← unaffiliated media falls back to type+month
  ├── Audio/2025-08/
  └── Other/2025-07/

Install dependencies (one-time):
  pip3 install pymupdf python-docx openpyxl python-pptx

Usage:
  python3 scripts/sort_documents.py ~/Documents
  python3 scripts/sort_documents.py ~/Documents --dry-run
  python3 scripts/sort_documents.py ~/Documents --output ~/Organized_Docs
  python3 scripts/sort_documents.py ~/Documents --copy              # keep originals
  python3 scripts/sort_documents.py ~/Documents --no-projects       # skip name grouping
  python3 scripts/sort_documents.py ~/Documents --vault Organized/Vault  # use Obsidian
"""

import argparse
import csv
import os
import re
import shutil
import sys
from collections import defaultdict
from datetime import datetime
from pathlib import Path

# ── Optional library imports ──────────────────────────────────────────────────

try:
    import fitz  # pymupdf
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ── Supported file types ──────────────────────────────────────────────────────

EXTENSIONS = {
    # Documents
    '.pdf':     'PDF',
    '.doc':     'Word',   '.docx':    'Word',
    '.xls':     'Excel',  '.xlsx':    'Excel',
    '.csv':     'Spreadsheet',
    '.ppt':     'PowerPoint', '.pptx': 'PowerPoint',
    '.txt':     'Text',   '.md':      'Text',    '.rtf':  'Text',
    '.pages':   'Pages',  '.numbers': 'Numbers', '.key':  'Keynote',
    '.epub':    'eBook',
    # Video
    '.mov':     'Video',  '.mp4':     'Video',   '.m4v':  'Video',
    '.avi':     'Video',  '.mkv':     'Video',   '.wmv':  'Video',
    '.webm':    'Video',  '.flv':     'Video',   '.3gp':  'Video',
    '.mts':     'Video',  '.prores':  'Video',
    # Audio
    '.mp3':     'Audio',  '.m4a':     'Audio',   '.wav':  'Audio',
    '.aac':     'Audio',  '.flac':    'Audio',   '.ogg':  'Audio',
    '.aiff':    'Audio',  '.wma':     'Audio',   '.opus': 'Audio',
    # Design files
    '.psd':     'Design', '.ai':      'Design',  '.sketch': 'Design',
    '.fig':     'Design', '.indd':    'Design',  '.svg':  'Design',
    '.eps':     'Design', '.tif':     'Design',  '.tiff': 'Design',
    '.raw':     'Design', '.cr2':     'Design',  '.nef':  'Design',
    '.dng':     'Design',
    # Images (project assets — camera roll goes through quick_sort.py)
    '.jpg':     'Image',  '.jpeg':    'Image',   '.png':  'Image',
    '.heic':    'Image',  '.heif':    'Image',   '.gif':  'Image',
    '.webp':    'Image',  '.bmp':     'Image',
    # Web / code
    '.html':    'Web',    '.htm':     'Web',     '.css':  'Web',
    '.js':      'Web',    '.json':    'Web',     '.xml':  'Web',
    # Fonts
    '.otf':     'Font',   '.ttf':     'Font',    '.woff': 'Font',
    '.woff2':   'Font',
    # Archives / packages
    '.zip':     'Archive', '.gz':     'Archive', '.tar':  'Archive',
    '.rar':     'Archive', '.7z':     'Archive', '.dmg':  'Archive',
    '.sit':     'Archive', '.cpgz':   'Archive',
    # Legacy / misc formats
    '.bin':     'Misc',   '.mbox':    'Misc',    '.db':   'Misc',
    '.plist':   'Misc',   '.webloc':  'Misc',    '.ico':  'Misc',
    '.eot':     'Font',   '.afm':     'Font',    '.pfb':  'Font',
    '.pfm':     'Font',
    '.xlsm':    'Excel',  '.dot':     'Word',    '.odt':  'Word',
    '.thmx':    'Design', '.merlin2': 'Misc',
    '.liquid':  'Web',    '.key':     'Keynote',
    '.textClipping': 'Misc', '.url':  'Misc',
}

# Types that are media/binary (no text extraction — classification relies
# entirely on filename, folder path, and vault entity matching)
MEDIA_TYPES = {'Video', 'Audio', 'Design', 'Image', 'Archive', 'Web', 'Font', 'Misc'}

# ── Category definitions ──────────────────────────────────────────────────────
# Each entry: category_name → (keywords_list, weight)
# Weight amplifies the score so high-confidence categories win ties.

CATEGORIES = {
    'Finance': (
        ['invoice', 'receipt', 'payment', 'bank statement', 'account statement',
         'tax return', 'w-2', 'w2', '1099', '1040', 'balance sheet', 'transaction',
         'deposit', 'wire transfer', 'irs', 'refund', 'expense report', 'payroll',
         'salary', 'loan', 'mortgage', 'credit card', 'brokerage', 'dividend',
         'capital gains', 'amount due', 'total due', 'billing', 'estimate',
         'purchase order', 'accounts payable', 'accounts receivable'],
        2.0
    ),
    'Legal': (
        ['agreement', 'contract', 'lease agreement', 'terms and conditions',
         'privacy policy', 'non-disclosure', 'nda', 'liability', 'attorney',
         'court', 'lawsuit', 'plaintiff', 'defendant', 'notarized', 'amendment',
         'whereas', 'indemnify', 'arbitration', 'jurisdiction', 'governing law',
         'intellectual property', 'trademark', 'copyright', 'patent', 'herein',
         'licensee', 'licensor', 'effective date', 'binding agreement'],
        2.0
    ),
    'Medical': (
        ['patient', 'diagnosis', 'prescription', 'physician', 'hospital', 'clinic',
         'insurance claim', 'medical record', 'pharmacy', 'lab results',
         'blood work', 'treatment plan', 'referral', 'explanation of benefits',
         'eob', 'deductible', 'copay', 'hsa', 'fsa', 'provider', 'dosage',
         'medication', 'symptoms', 'procedure', 'date of service', 'icd',
         'cpt code', 'primary care', 'specialist'],
        2.0
    ),
    'Real_Estate': (
        ['property', 'listing', 'appraisal', 'escrow', 'title report', 'deed',
         'hoa', 'square feet', 'zillow', 'realtor', 'closing costs',
         'purchase agreement', 'home inspection', 'refinance', 'loan estimate',
         'disclosure', 'mls', 'offer letter', 'earnest money', 'contingency'],
        2.0
    ),
    'Work': (
        ['proposal', 'scope of work', 'statement of work', 'meeting notes', 'agenda',
         'deliverable', 'client', 'timeline', 'stakeholder', 'quarterly review',
         'annual report', 'strategy', 'roadmap', 'milestone', 'kpi', 'okr',
         'performance review', 'forecast', 'executive summary', 'action items',
         'project plan', 'status update', 'go-to-market'],
        1.5
    ),
    'Personal': (
        ['journal', 'diary', 'personal note', 'dear ', 'sincerely', 'memoir',
         'birthday', 'anniversary', 'family', 'vacation', 'wish list', 'goals',
         'new year', 'gratitude', 'reflection', 'bucket list', 'love letter'],
        1.0
    ),
    'Reference': (
        ['user manual', 'user guide', 'tutorial', 'how to', 'instructions',
         'documentation', 'reference guide', 'handbook', 'specification',
         'glossary', 'faq', 'readme', 'getting started', 'table of contents',
         'appendix', 'index', 'quick start'],
        1.0
    ),
    'Creative': (
        ['screenplay', 'chapter one', 'chapter 1', 'fade in', 'ext.', 'int.',
         'scene', 'act one', 'act 1', 'protagonist', 'antagonist', 'lyrics',
         'verse', 'chorus', 'manuscript', 'short story', 'plot outline',
         'character bio', 'logline'],
        1.5
    ),
}

# ── Vault entity loading ──────────────────────────────────────────────────────

# Single-word entries that are almost certainly OCR noise, not real names
VAULT_NOISE = {
    'ahhh', 'anyone', 'another', 'autistic', 'beautiful', 'because', 'call',
    'captured', 'chat', 'church', 'clock', 'code', 'camp', 'boss', 'boomerang',
    'amazing', 'approved', 'arriving', 'abandoned', 'art', 'animation',
    'add', 'al', 'around', 'also', 'always', 'already', 'after',
    'acne', 'amma', 'alto', 'beane', 'cath', 'chuck', 'coliv',
}

# Known company/service suffixes for classification
COMPANY_SUFFIXES = {
    'inc', 'llc', 'corp', 'ltd', 'co', 'company', 'group', 'studio', 'studios',
    'agency', 'partners', 'consulting', 'services', 'solutions', 'technologies',
    'tech', 'media', 'capital', 'ventures', 'fund', 'labs', 'lab',
}


def _parse_entity_file(md_path):
    """Parse a vault entity markdown file → list of entity name strings."""
    entities = []
    try:
        text = md_path.read_text(encoding='utf-8', errors='ignore')
        # Match lines like: - **Entity Name** — [[...]]
        for m in re.finditer(r'^- \*\*(.+?)\*\*', text, re.MULTILINE):
            name = m.group(1).strip()
            # Clean up OCR multi-line artifacts (newlines inside the name)
            name = re.sub(r'\s+', ' ', name)
            entities.append(name)
    except Exception:
        pass
    return entities


def _is_quality_entity(name):
    """Filter out OCR noise — keep real people/company names."""
    if len(name) < 4:
        return False
    words = name.split()
    lower = name.lower()
    # All-caps short strings are usually OCR noise (APPROVED, AMALGACLEAR)
    if name.isupper() and len(words) == 1:
        return False
    # Single word in noise list
    if len(words) == 1 and lower in VAULT_NOISE:
        return False
    # Single generic word (verb, article, common noun)
    if len(words) == 1 and lower in {
        'add', 'and', 'art', 'call', 'chat', 'code', 'done', 'from', 'get',
        'got', 'has', 'its', 'let', 'new', 'not', 'now', 'old', 'one',
        'our', 'out', 'put', 'set', 'the', 'this', 'top', 'use', 'was',
        'who', 'will', 'with', 'you',
    }:
        return False
    return True


def _names_from_filenames(directory, prefix=''):
    """Extract entity names from markdown filenames in a directory.
    Strips optional prefix (e.g. 'PRJ - ') and .md extension."""
    names = set()
    d = Path(directory)
    if not d.is_dir():
        return names
    for f in d.iterdir():
        if f.suffix.lower() != '.md':
            continue
        name = f.stem
        if prefix and name.startswith(prefix):
            name = name[len(prefix):]
        name = name.strip(' -–—')
        if name and len(name) >= 3:
            names.add(name)
    return names


def load_vault_entities(vault_dir):
    """
    Load known contacts and companies from the Obsidian vault entity indexes.
    Returns two sets: (contacts, companies)
      contacts  — real people names (e.g. "Charlie Smith", "Andrea Vanzo")
      companies — businesses/services (e.g. "Your Company", "Apple Music")
    """
    vault = Path(vault_dir)
    contacts  = set()
    companies = set()

    entity_dir = vault / 'Entities'
    if not entity_dir.exists():
        return contacts, companies

    # Contacts → People/ folder
    contacts_file = entity_dir / 'Contacts.md'
    if contacts_file.exists():
        for name in _parse_entity_file(contacts_file):
            if _is_quality_entity(name):
                if len(name.split()) >= 2:
                    contacts.add(name)

    # Services → Companies/ folder
    services_file = entity_dir / 'Services.md'
    if services_file.exists():
        for name in _parse_entity_file(services_file):
            if _is_quality_entity(name):
                words = name.split()
                lower_words = [w.lower() for w in words]
                if len(words) >= 2 or any(w in COMPANY_SUFFIXES for w in lower_words):
                    companies.add(name)

    return contacts, companies


def load_vault_dir(vault_dir_dir):
    """
    Load project and company names from the YOURBRAND vault.
    Returns a dict: { lowercase_name: canonical_name }
    Used for folder routing — if a file's path or content mentions a known
    project/company, it goes into a folder named after that entity.
    """
    vault = Path(vault_dir_dir)
    known = {}  # lowercase → canonical display name

    # 05_PROJECTS — Active, Archived, Dormant, Incubating
    projects_dir = vault / '05_PROJECTS'
    if projects_dir.is_dir():
        for status_dir in projects_dir.iterdir():
            if status_dir.is_dir():
                for name in _names_from_filenames(status_dir, prefix='PRJ - '):
                    # Skip meta notes (YOURBRAND Studio - ...)
                    if name.startswith('YOURBRAND Studio'):
                        continue
                    known[name.lower()] = name

    # 06_OUTPUTS/Business Network — companies, collaborators, orgs
    biz_dir = vault / '06_OUTPUTS' / 'Business Network'
    if biz_dir.is_dir():
        for name in _names_from_filenames(biz_dir):
            known[name.lower()] = name

    # 04_CANON/Studio/Companies
    companies_dir = vault / '04_CANON' / 'Studio' / 'Companies'
    if companies_dir.is_dir():
        for name in _names_from_filenames(companies_dir):
            known[name.lower()] = name
        # Also check subfolders (e.g. Your Company Collection/)
        for sub in companies_dir.iterdir():
            if sub.is_dir():
                known[sub.name.lower()] = sub.name

    # 04_CANON/YOURBRAND/Artworks — artwork names as projects
    artworks_dir = vault / '04_CANON' / 'YOURBRAND' / 'Artworks'
    if artworks_dir.is_dir():
        for f in artworks_dir.iterdir():
            if f.suffix.lower() == '.md':
                # Format: "2016 - Cash Bricks ($100,000).md"
                name = f.stem
                # Strip year prefix
                m = re.match(r'^\d{4}\s*[-–—]\s*(.+)$', name)
                if m:
                    name = m.group(1).strip()
                if name and len(name) >= 4:
                    known[name.lower()] = name

    return known


def match_vault_entity(text, filename, folder_hint, vault_known):
    """
    Check if text/filename/folder matches a known YOURBRAND project or company.
    Returns canonical name or None.
    Tries longest matches first to avoid false positives (e.g. "Art" matching
    when the real entity is "Your Company").
    """
    if not vault_known:
        return None

    combined = (folder_hint + ' ' + filename + ' ' + text[:3000]).lower()

    # Sort by length descending — match "Your Company" before "Art"
    for key in sorted(vault_known.keys(), key=len, reverse=True):
        if len(key) < 4:
            continue
        if key in combined:
            return vault_known[key]

    return None


def match_vault_entity(text, filename, contacts, companies):
    """
    Return (folder_type, entity_name) if a known vault entity is found.
    folder_type is 'People' or 'Companies', or None if no match.
    """
    combined = (filename + ' ' + text[:2000]).lower()

    # Check contacts first (more specific)
    for name in contacts:
        if name.lower() in combined:
            safe = re.sub(r'[^\w\s-]', '', name).strip().replace(' ', '_')
            return 'People', safe

    # Then companies
    for name in companies:
        if name.lower() in combined:
            safe = re.sub(r'[^\w\s-]', '', name).strip().replace(' ', '_')
            return 'Companies', safe

    return None, None


# Tokens to ignore when detecting project names from filenames
IGNORED_TOKENS = {
    'the', 'and', 'for', 'from', 'with', 'this', 'that', 'copy', 'draft',
    'final', 'new', 'old', 'backup', 'archive', 'temp', 'test', 'doc',
    'pdf', 'file', 'document', 'untitled', 'updated', 'revised', 'signed',
    'completed', 'done', 'original', 'version', 'rev', 'scan', 'scanned',
    'img', 'image', 'photo', 'misc', 'other', 'general', 'notes',
}

# ── Text extraction ───────────────────────────────────────────────────────────

def extract_text_pdf(path):
    if not HAS_PYMUPDF:
        return ''
    try:
        doc = fitz.open(str(path))
        pages = []
        for i, page in enumerate(doc):
            if i >= 5:  # first 5 pages is enough to classify
                break
            pages.append(page.get_text())
        doc.close()
        return ' '.join(pages)
    except Exception:
        return ''


def extract_text_docx(path):
    if not HAS_DOCX:
        return ''
    try:
        doc = DocxDocument(str(path))
        return ' '.join(p.text for p in doc.paragraphs[:150])
    except Exception:
        return ''


def extract_text_xlsx(path):
    if not HAS_OPENPYXL:
        return ''
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        text = []
        for sheet in wb.worksheets[:3]:      # first 3 sheets
            for row in sheet.iter_rows(max_row=60, values_only=True):
                for cell in row:
                    if cell and isinstance(cell, str):
                        text.append(cell)
        wb.close()
        return ' '.join(text)
    except Exception:
        return ''


def extract_text_pptx(path):
    if not HAS_PPTX:
        return ''
    try:
        prs = Presentation(str(path))
        text = []
        for slide in prs.slides[:10]:       # first 10 slides
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        return ' '.join(text)
    except Exception:
        return ''


def extract_text_csv(path):
    try:
        rows = []
        with open(path, newline='', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i >= 30:
                    break
                rows.append(' '.join(str(c) for c in row))
        return ' '.join(rows)
    except Exception:
        return ''


def extract_text_plain(path):
    try:
        return path.read_text(encoding='utf-8', errors='ignore')[:8000]
    except Exception:
        return ''


def extract_text(path):
    """Dispatch to the right extractor based on file extension."""
    ext = path.suffix.lower()
    if ext == '.pdf':
        return extract_text_pdf(path)
    elif ext in ('.doc', '.docx'):
        return extract_text_docx(path)
    elif ext in ('.xls', '.xlsx'):
        return extract_text_xlsx(path)
    elif ext in ('.ppt', '.pptx'):
        return extract_text_pptx(path)
    elif ext == '.csv':
        return extract_text_csv(path)
    elif ext in ('.txt', '.md', '.rtf'):
        return extract_text_plain(path)
    return ''


# ── Classification ────────────────────────────────────────────────────────────

def classify(text, filename):
    """Return (category, score). Score 0 = unclassified → use type fallback."""
    combined = (filename + ' ' + text).lower()
    scores = {}
    for cat, (keywords, weight) in CATEGORIES.items():
        hits = sum(1 for kw in keywords if kw in combined)
        if hits:
            scores[cat] = hits * weight
    if not scores:
        return 'Other', 0.0
    best = max(scores, key=scores.get)
    return best, scores[best]


# ── Date extraction ───────────────────────────────────────────────────────────

def get_date(path):
    """Try filename date patterns first, then fall back to file birth time."""
    name = path.stem
    # Full date: YYYY-MM-DD / YYYY_MM_DD / YYYYMMDD
    m = re.search(r'(20\d{2})[-_.]?(0[1-9]|1[0-2])[-_.]?(0[1-9]|[12]\d|3[01])', name)
    if m:
        try:
            return datetime(int(m.group(1)), int(m.group(2)), int(m.group(3)))
        except ValueError:
            pass
    # Month only: YYYY-MM
    m = re.search(r'(20\d{2})[-_.](0[1-9]|1[0-2])', name)
    if m:
        return datetime(int(m.group(1)), int(m.group(2)), 1)
    # File birth time (macOS) or modification time
    st = os.stat(path)
    ts = getattr(st, 'st_birthtime', None) or st.st_mtime
    return datetime.fromtimestamp(ts)


# ── Project name detection ────────────────────────────────────────────────────

def tokenize(name):
    """Break a filename stem into meaningful project-name tokens."""
    # Strip date patterns
    name = re.sub(r'20\d{2}[-_.]?\d{2}[-_.]?\d{0,2}', '', name)
    # Strip version patterns (v1, v2.1, rev3)
    name = re.sub(r'\b(v\d+(\.\d+)?|rev\d+)\b', '', name, flags=re.IGNORECASE)
    # Split on anything non-alphanumeric
    tokens = re.split(r'[^a-zA-Z0-9]+', name)
    return [t.lower() for t in tokens if len(t) >= 4 and t.lower() not in IGNORED_TOKENS]


def detect_projects(files, source_root):
    """
    Cluster files into named projects based on shared filename tokens.
    A token qualifies as a project name if it appears in 3–12 files
    AND those files span at least 2 different top-level source folders
    (proving it's a real project name, not just a batch of similarly-named exports).
    Returns: dict of filepath → project_name (title-cased), or empty string.
    """
    token_to_files = defaultdict(list)
    for f in files:
        for tok in tokenize(f.stem):
            token_to_files[tok].append(f)

    # Keep only tokens that look like real project names:
    # - appear in 3+ files (not just a coincidental pair)
    # - appear in no more than 12 (not a generic word)
    # - are purely alphabetic (no hashes, IDs, dimensions like 42X30)
    # - span 2+ different TOP-LEVEL source folders
    project_tokens = {}
    for tok, flist in token_to_files.items():
        if not (3 <= len(flist) <= 12):
            continue
        if not tok.isalpha():
            continue
        if len(tok) < 5:
            continue
        # Check that files come from 2+ distinct top-level folders
        top_folders = set()
        for f in flist:
            try:
                rel = f.relative_to(source_root)
                top_folders.add(rel.parts[0] if len(rel.parts) > 1 else '')
            except ValueError:
                top_folders.add('')
        if len(top_folders) < 2:
            continue
        project_tokens[tok] = flist

    file_project = {}
    for f in files:
        tokens = tokenize(f.stem)
        best_tok, best_count = None, 0
        for t in tokens:
            if t in project_tokens and len(project_tokens[t]) > best_count:
                best_tok = t
                best_count = len(project_tokens[t])
        file_project[f] = best_tok.title() if best_tok else ''

    return file_project


# ── Source folder path analysis ────────────────────────────────────────────────

# Generic folder names that carry no meaning for classification
GENERIC_FOLDERS = {
    'documents', 'document', 'docs', 'files', 'file', 'downloads', 'download',
    'desktop', 'users', 'user', 'home', 'volumes', 'tmp', 'temp',
    'library', 'data', 'content', 'stuff', 'misc', 'old', 'new',
    'backup', 'backups', 'archive', 'archives', 'shared',
    'inbox', 'general', 'media', 'proxy', 'assets', 'horizontal', 'vertical',
    'audio', 'video', 'videos', 'images', 'photos', 'pictures',
    'decks and docs', 'file cabinet', 'file misc', 'forms', 'resources',
    'financial models', 'contact lists', 'inventory lists',
    'chatgpt download', 'gmail downloads', 'granola downloads',
}


def get_folder_context(filepath, source_root):
    """
    Extract meaningful folder names from the path between source_root and the file.

    Example: source is ~/Documents/MyDocs
             file is   ~/Documents/MyDocs/Your Company/Contracts/2024/lease.pdf
             returns:  ['Your Company', 'Contracts', '2024']

    Returns a list of folder name strings, filtering out generic names.
    """
    try:
        rel = filepath.parent.relative_to(source_root)
    except ValueError:
        return []

    parts = [p for p in rel.parts if p.lower() not in GENERIC_FOLDERS]
    return parts


def get_subpath_below_project(filepath, source_root, project_folder_name):
    """
    Given a file and the project folder name it matched on, return the
    relative subfolder path BELOW that project folder.

    Example: source_root = /new
             filepath    = /new/CF Holdings/Clients/Alfred Coffee/CBS/Mockups/file.png
             project_folder_name matched on = 'CF Holdings'
             returns:     Path('Clients/Alfred Coffee/CBS/Mockups')

    This preserves internal folder structure within a project.
    Returns Path('') if the file is directly in the project root.
    """
    try:
        rel = filepath.parent.relative_to(source_root)
    except ValueError:
        return Path('')

    parts = rel.parts
    # Find where the project folder name appears in the path
    proj_lower = project_folder_name.lower()
    for i, part in enumerate(parts):
        if part.lower() == proj_lower or proj_lower in part.lower():
            # Everything after this folder is the internal subpath
            sub_parts = parts[i + 1:]
            return Path(*sub_parts) if sub_parts else Path('')
    # If folder name doesn't appear literally (e.g. YOURBRAND vault matched on content
    # within "Your Company Videos"), try to find the first meaningful folder and keep below it
    for i, part in enumerate(parts):
        if part.lower() not in GENERIC_FOLDERS:
            sub_parts = parts[i + 1:]
            return Path(*sub_parts) if sub_parts else Path('')
    return Path('')


def folder_to_project(folder_parts):
    """
    If the source folder path contains a meaningful project/entity name,
    return it as a cleaned project name. Otherwise return ''.

    Heuristic: the first non-date, non-generic, non-UUID folder name is the project.
    """
    for part in folder_parts:
        # Skip pure date folders like "2024" or "2024-01"
        if re.match(r'^20\d{2}(-\d{2})?$', part):
            continue
        # Skip single-char or very short
        if len(part) < 3:
            continue
        # Skip UUIDs and hex hashes (ChatGPT export IDs, etc.)
        if re.match(r'^[0-9a-f]{8}(-[0-9a-f]{4}){3}-[0-9a-f]{12}$', part, re.IGNORECASE):
            continue
        if re.match(r'^(user-)?[A-Za-z0-9]{20,}$', part):
            continue
        # Clean up for use as a folder name
        safe = re.sub(r'[^\w\s-]', '', part).strip()
        if safe:
            return safe
    return ''


# ── Destination resolution ────────────────────────────────────────────────────

def unique_dest(dest_dir, filename):
    """Return a collision-free destination path, respecting macOS 255-byte limit."""
    stem = Path(filename).stem
    suffix = Path(filename).suffix
    # Truncate if filename is too long (macOS limit: 255 bytes)
    max_stem = 255 - len(suffix.encode('utf-8')) - 5  # reserve room for _999
    if len(stem.encode('utf-8')) > max_stem:
        while len(stem.encode('utf-8')) > max_stem:
            stem = stem[:-1]
        filename = stem + suffix
    dest = dest_dir / filename
    if not dest.exists():
        return dest
    counter = 1
    while dest.exists():
        dest = dest_dir / f'{stem}_{counter}{suffix}'
        counter += 1
    return dest


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description='Organize documents by content and name',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__
    )
    parser.add_argument('source_dir',
                        help='Directory to scan for documents')
    parser.add_argument('--output', '-o',
                        help='Output directory (default: <source_dir>/../Organized_Docs)')
    parser.add_argument('--dry-run', '-n', action='store_true',
                        help='Preview without touching files')
    parser.add_argument('--copy', action='store_true',
                        help='Copy instead of move (keeps originals in place)')
    parser.add_argument('--no-projects', action='store_true',
                        help='Skip project name grouping')
    parser.add_argument('--min-score', type=float, default=2.0,
                        help='Minimum classification score to use a category (default: 2.0)')
    parser.add_argument('--vault', '-v',
                        help='Path to screenshot Obsidian vault (loads contacts & companies)')
    parser.add_argument('--yourbrand-vault',
                        help='Path to YOURBRAND vault (loads projects, companies, artworks)')
    parser.add_argument('--recursive', action='store_true', default=True,
                        help='Scan subdirectories (default: on)')
    args = parser.parse_args()

    source = Path(args.source_dir).expanduser().resolve()
    if not source.is_dir():
        print(f'Error: {source} is not a directory')
        sys.exit(1)

    output = (
        Path(args.output).expanduser().resolve()
        if args.output
        else source.parent / 'Organized_Docs'
    )

    # Check library availability
    missing = []
    if not HAS_PYMUPDF:   missing.append('pymupdf (PDFs)')
    if not HAS_DOCX:      missing.append('python-docx (Word)')
    if not HAS_OPENPYXL:  missing.append('openpyxl (Excel)')
    if not HAS_PPTX:      missing.append('python-pptx (PowerPoint)')
    if missing:
        print('⚠  Missing libraries (install for full support):')
        print(f'   pip3 install {" ".join(p.split()[0] for p in missing)}')
        print()

    # Guard: don't let output live inside source (would re-process own output)
    try:
        output.relative_to(source)
        print(f'Error: output dir ({output}) is inside source dir ({source})')
        print(f'       Use --output to set a different location')
        sys.exit(1)
    except ValueError:
        pass  # good — they're separate

    print('=' * 50)
    print('Document Organizer')
    print('=' * 50)
    print(f'Source:  {source}')
    print(f'Output:  {output}')
    print(f'Mode:    {"DRY RUN" if args.dry_run else "COPY" if args.copy else "MOVE"}')
    if args.vault:
        print(f'Vault:   {args.vault}')
    print()

    # ── Collect files ──
    all_files = []
    scan_fn = source.rglob if args.recursive else source.glob
    for path in scan_fn('*'):
        if not path.is_file():
            continue
        if path.is_symlink():
            continue
        if path.name.startswith('.'):
            continue
        ext = path.suffix.lower()
        if ext not in EXTENSIONS and ext != '':
            # Unknown extension — still process if inside a project folder
            EXTENSIONS[ext] = 'Misc'
        if not ext:
            # Extensionless file — treat as Misc
            EXTENSIONS[''] = 'Misc'
        # Skip files already inside the output folder
        try:
            path.relative_to(output)
            continue
        except ValueError:
            pass
        all_files.append(path)

    total = len(all_files)
    if total == 0:
        print('No supported documents found.')
        return

    type_counts = defaultdict(int)
    for f in all_files:
        type_counts[EXTENSIONS[f.suffix.lower()]] += 1

    print(f'Found {total} files:')
    for doc_type, count in sorted(type_counts.items(), key=lambda x: -x[1]):
        print(f'  {doc_type:<15} {count:>5}')
    print()

    # ── Load vault entities ──
    vault_contacts, vault_companies = set(), set()
    if args.vault:
        vault_path = Path(args.vault).expanduser().resolve()
        if not (vault_path / 'Entities').is_dir():
            print(f'⚠  Vault path has no Entities/ folder: {vault_path}')
            print(f'   Continuing without vault matching.')
            print()
        else:
            vault_contacts, vault_companies = load_vault_entities(vault_path)
            print(f'Vault entities loaded:')
            print(f'  Contacts (people):  {len(vault_contacts)}')
            print(f'  Services/companies: {len(vault_companies)}')
            print()

    # ── Load YOURBRAND vault entities ──
    vault_known = {}
    if args.vault_dir:
        vault_path = Path(args.vault_dir).expanduser().resolve()
        if not vault_path.is_dir():
            print(f'⚠  YOURBRAND vault path not found: {vault_path}')
            print()
        else:
            vault_known = load_vault_dir(vault_path)
            print(f'YOURBRAND vault loaded: {len(vault_known)} known projects/companies/artworks')
            # Show a sample
            sample = sorted(vault_known.values())[:15]
            for name in sample:
                print(f'  • {name}')
            if len(vault_known) > 15:
                print(f'  ... and {len(vault_known) - 15} more')
            print()

    # ���─ Project detection ──
    file_project = {}
    if not args.no_projects:
        file_project = detect_projects(all_files, source)
        project_names = sorted(set(v for v in file_project.values() if v))
        if project_names:
            print(f'Detected {len(project_names)} project group(s):')
            for p in project_names:
                count = sum(1 for v in file_project.values() if v == p)
                print(f'  Projects/{p}/  ({count} files)')
            print()

    # ── Process each file ──
    placements = []   # list of (src_path, dest_path, label)
    stats = defaultdict(int)

    for i, filepath in enumerate(sorted(all_files), 1):
        if i % 25 == 0 or i == total:
            pct = int(i / total * 100)
            print(f'  Classifying: {i}/{total} ({pct}%)', end='\r', flush=True)

        file_date = get_date(filepath)
        month_str = file_date.strftime('%Y-%m')
        doc_type  = EXTENSIONS.get(filepath.suffix.lower(), 'Other')

        # Get source folder context (e.g. ["Your Company", "Contracts", "2024"])
        folder_parts = get_folder_context(filepath, source)
        folder_hint  = ' '.join(folder_parts)  # flat string for keyword matching

        # Step 0: source folder has a meaningful project/entity name
        folder_project = folder_to_project(folder_parts)

        # ── Step 0: YOURBRAND vault match on folder name ──
        # Check if folder name IS a known project (highest confidence signal)
        vault_match = match_vault_entity('', '', folder_hint, vault_known) if folder_hint else None

        project = file_project.get(filepath, '')
        if vault_match:
            safe_name = vault_match.replace(' ', '_')
            subpath = get_subpath_below_project(filepath, source, vault_match)
            dest_dir = output / 'Projects' / safe_name / subpath
            label    = f'Projects/{safe_name}'
        elif folder_project and not project:
            folder_type, entity_name = match_vault_entity(
                '', folder_project, vault_contacts, vault_companies
            )
            if entity_name:
                subpath = get_subpath_below_project(filepath, source, folder_project)
                dest_dir = output / folder_type / entity_name / subpath
                label    = f'{folder_type}/{entity_name}'
            else:
                safe_name = folder_project.replace(' ', '_')
                subpath = get_subpath_below_project(filepath, source, folder_project)
                dest_dir = output / 'Projects' / safe_name / subpath
                label    = f'Projects/{safe_name}'
        elif project:
            dest_dir = output / 'Projects' / project
            label    = f'Projects/{project}'
        else:
            is_media = doc_type in MEDIA_TYPES
            text = '' if is_media else extract_text(filepath)
            combined_text = folder_hint + ' ' + text if folder_hint else text

            vault_content_match = match_vault_entity(
                combined_text, filepath.stem, folder_hint, vault_known
            )
            if vault_content_match:
                safe_name = vault_content_match.replace(' ', '_')
                subpath = get_subpath_below_project(filepath, source, '')
                dest_dir = output / 'Projects' / safe_name / subpath
                label    = f'Projects/{safe_name}'
            else:
                folder_type, entity_name = match_vault_entity(
                    combined_text, filepath.stem, vault_contacts, vault_companies
                )
                if entity_name:
                    dest_dir = output / folder_type / entity_name
                    label    = f'{folder_type}/{entity_name}'
                elif not is_media:
                    category, score = classify(combined_text, filepath.stem)
                    if score >= args.min_score:
                        dest_dir = output / category / month_str
                        label    = f'{category}/{month_str}'
                    else:
                        dest_dir = output / doc_type / month_str
                        label    = f'{doc_type}/{month_str}'
                else:
                    dest_dir = output / doc_type / month_str
                    label    = f'{doc_type}/{month_str}'

        dest = unique_dest(dest_dir, filepath.name)
        placements.append((filepath, dest, label))
        stats[label.split('/')[0]] += 1

    print(f'  Classifying: {total}/{total} (100%) — done          ')
    print()

    # ── Preview or apply ──
    if args.dry_run:
        print(f'PREVIEW — first 30 of {total} files:')
        print()
        for src, dst, label in placements[:30]:
            rel = dst.relative_to(output)
            # Show source folder context if file isn't at root of source
            src_ctx = get_folder_context(src, source)
            src_display = '/'.join(src_ctx + [src.name]) if src_ctx else src.name
            print(f'  [{label.split("/")[0]}] {src_display}')
            print(f'       → {rel}')
        if total > 30:
            print(f'\n  ... and {total - 30} more (run without --dry-run to apply)')
    else:
        action_fn = shutil.copy2 if args.copy else shutil.move
        action_name = 'Copying' if args.copy else 'Moving'
        errors = []
        manifest = []
        for j, (src, dst, lbl) in enumerate(placements, 1):
            if j % 25 == 0 or j == total:
                print(f'  {action_name}: {j}/{total}', end='\r', flush=True)
            try:
                dst.parent.mkdir(parents=True, exist_ok=True)
                action_fn(str(src), str(dst))
                manifest.append((str(src), str(dst)))
            except Exception as e:
                errors.append((str(src), str(e)))
        print(f'  {action_name}: {total}/{total} — done          ')

        # Write manifest so moves can be reversed
        manifest_path = output / '_manifest.csv'
        with open(manifest_path, 'w', newline='') as mf:
            writer = csv.writer(mf)
            writer.writerow(['source', 'destination'])
            writer.writerows(manifest)
        print(f'\n  Manifest saved: {manifest_path}')

        if errors:
            print(f'\n  ⚠  {len(errors)} file(s) failed:')
            for path, err in errors[:10]:
                print(f'     {Path(path).name}: {err}')
            if len(errors) > 10:
                print(f'     ... and {len(errors) - 10} more')

    # ── Summary ──
    print()
    print('Results by category:')
    for cat, count in sorted(stats.items(), key=lambda x: -x[1]):
        bar = '█' * min(count, 30)
        print(f'  {cat:<20} {count:>5}  {bar}')
    print(f'  {"─" * 40}')
    print(f'  {"Total":<20} {total:>5}')

    if not args.dry_run:
        print()
        print(f'Output: {output}')
        try:
            import subprocess
            result = subprocess.run(['du', '-sh', str(output)], capture_output=True, text=True)
            if result.stdout:
                print(f'Size:   {result.stdout.split()[0]}')
        except Exception:
            pass


if __name__ == '__main__':
    main()
